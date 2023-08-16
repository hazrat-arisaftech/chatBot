import os
from fastapi import FastAPI, File, UploadFile, Form 
from typing import Annotated
from pydantic import BaseModel
import PyPDF2
from docx import Document
from pptx import Presentation
import tempfile
import logging
from services import EchoService
from pymongo import MongoClient
from fastapi.responses import JSONResponse, FileResponse
from bson.json_util import dumps
from bson.objectid import ObjectId 
from datetime import datetime
from pathlib import Path
# from pymongo.errors import InvalidId

# from bson import ObjectId 
# import Request

logging.config.fileConfig('logging.conf', disable_existing_loggers=False)
logger = logging.getLogger(__name__)

# class Prompt(BaseModel):
 

app = FastAPI()

client = MongoClient("mongodb://localhost:27017")
# db = client("file_storage") 
db = client.file_storage
file_collection = db["files"]

class FileUpload(BaseModel):
    name: str
    content: bytes
    time: str
    
    
@app.get("/health")
async def health():
    # logger.info("logging from the root logger")
    # EchoService.echo("hi")
    return {"Message": "Up and Running"}


@app.post("/upload/")
async def upload(prompt: Annotated[str, Form()], file: Annotated[UploadFile, File()]):
    temp_dir = tempfile.mkdtemp()
    file_path = os.path.join(temp_dir, file.filename)

    with open(file_path, "wb") as f:
        f.write(file.file.read())

    if (file.filename.endswith(".pdf")):
        print("It's a pdf file")
        reader = PyPDF2.PdfReader(file.file)
        for page in reader.pages:
            content = page.extract_text()
            print(content)

    elif (file.filename.endswith(".docx") or file.filename.endswith(".doc")):
        print("It's doc file")
        doc = Document(file_path)
        content = ""
        for para in doc.paragraphs:
            content += para.text
        print(content)
            
    elif (file.filename.endswith(".ppt") or file.filename.endswith(".pptx")):
        print("It's a ppt file")
        content = ""
        presentation = Presentation(file_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                if hasattr(shape, 'text'):
                    content +=  shape.text + "\n"
        print(content)  

    elif (file.filename.endswith(".txt")):
        print("It's a txt file")
        content = ""
        with open(file_path, 'r', encoding='utf-8') as f:
           content = f.read() 
        print(content)

    return {"prompt": prompt, "file": file.filename}

@app.post("/uploadfile")
async def upload_file(file: UploadFile = File(...)):

    current_datetime = datetime.now()
    
    formatted_datetime = current_datetime.strftime("%Y-%m-%d_%H-%M-%S")
    content = await file.read()
    file_data = {
        "name": file.filename,
        "content": content,
        "time": formatted_datetime
    }

    file_collection.insert_one(file_data)
    # DIR = "temp"
    if not os.path.exists("temp"):
        os.makedirs("temp")

    print(file_data["time"])


    with open(os.path.join("temp", file_data["time"]), "wb") as f:
        f.write(content)

    return {"Message": "File upload successful"}

@app.get("/all")
async def get_files():
    files = list (file_collection.find())
    serialized_files = dumps(files)
    return JSONResponse(content=serialized_files)

@app.get("/download/{file_id}")
async def download_file(file_id: str):
    
    file_data = file_collection.find_one({"_id": ObjectId(file_id)})

    if not file_data:
        return JSONResponse(content={"message": "File not found"}, status_code=404)
    
    file_content = file_data["content"]
    file_name = file_data["name"] 
    file_time = file_data["time"]

    if not os.path.exists("temp"):
        os.makedirs("temp")
    
    file_path = Path("temp") / file_time
    print(file_path)
    # return FileResponse(file_path, file_name)
    return FileResponse(file_path, media_type='application/octet-stream',filename=file_name)
    
    